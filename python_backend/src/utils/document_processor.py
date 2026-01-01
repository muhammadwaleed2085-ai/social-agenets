"""
Document Processing Utility

Parses PDF, DOCX, PPTX, XLSX, and other document formats into text for LLM consumption.
Uses lightweight libraries: pypdf, python-docx, python-pptx, openpyxl.
"""

import base64
import io
import logging
from typing import Optional

logger = logging.getLogger(__name__)

# Supported MIME types
SUPPORTED_DOCUMENT_TYPES = {
    # PDF
    'application/pdf': '.pdf',
    # Word documents
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
    'application/msword': '.doc',
    # PowerPoint
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
    'application/vnd.ms-powerpoint': '.ppt',
    # Excel
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
    'application/vnd.ms-excel': '.xls',
    # Text formats
    'text/plain': '.txt',
    'text/markdown': '.md',
    'text/html': '.html',
    'text/csv': '.csv',
    # Data formats
    'application/json': '.json',
    'application/xml': '.xml',
    'text/xml': '.xml',
}


def process_document_from_base64(
    data: str,
    mime_type: str,
    filename: Optional[str] = None
) -> str:
    """
    Process a document from base64-encoded data and extract text.
    
    Supports: PDF, DOCX, PPTX, XLSX, HTML, TXT, CSV, JSON, XML.
    
    Args:
        data: Base64-encoded document content
        mime_type: MIME type of the document
        filename: Optional filename for the document
        
    Returns:
        Extracted text content as markdown string
    """
    try:
        # Decode base64 data
        binary_data = base64.b64decode(data)
        ext = SUPPORTED_DOCUMENT_TYPES.get(mime_type, '.txt')
        doc_name = filename or f"document{ext}"
        
        logger.info(f"Processing document: {doc_name} ({mime_type})")
        
        # Route to appropriate processor based on type
        if mime_type == 'application/pdf':
            content = _process_pdf(binary_data, doc_name)
        elif mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
            content = _process_docx(binary_data, doc_name)
        elif mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
            content = _process_pptx(binary_data, doc_name)
        elif mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel']:
            content = _process_xlsx(binary_data, doc_name)
        elif mime_type == 'text/html':
            content = _process_html(binary_data, doc_name)
        elif mime_type == 'text/csv':
            content = _process_csv(binary_data, doc_name)
        elif mime_type in ['text/plain', 'text/markdown']:
            content = _process_text(binary_data, doc_name)
        elif mime_type == 'application/json':
            content = _process_json(binary_data, doc_name)
        elif mime_type in ['application/xml', 'text/xml']:
            content = _process_xml(binary_data, doc_name)
        else:
            # Unknown type - try to decode as text
            content = _process_text(binary_data, doc_name)
        
        return content
        
    except Exception as e:
        logger.error(f"Document processing error: {e}", exc_info=True)
        return f"[Document processing failed for '{filename or 'document'}': {str(e)}]"


def _process_pdf(binary_data: bytes, doc_name: str) -> str:
    """Process PDF files using pypdf."""
    try:
        from pypdf import PdfReader
        
        pdf_file = io.BytesIO(binary_data)
        reader = PdfReader(pdf_file)
        
        text_parts = [f"## Document: {doc_name}\n"]
        
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(f"### Page {page_num}\n\n{page_text.strip()}")
        
        if len(text_parts) == 1:
            return f"[No text could be extracted from PDF: {doc_name}]"
        
        result = "\n\n".join(text_parts)
        logger.info(f"PDF processed: {doc_name} - {len(result)} chars extracted")
        return result
        
    except ImportError:
        return f"[PDF support requires pypdf. Document: {doc_name}]"
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        return f"[PDF extraction failed for '{doc_name}': {str(e)}]"


def _process_docx(binary_data: bytes, doc_name: str) -> str:
    """Process DOCX files using python-docx."""
    try:
        from docx import Document
        
        docx_file = io.BytesIO(binary_data)
        doc = Document(docx_file)
        
        text_parts = [f"## Document: {doc_name}\n"]
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                # Check if it's a heading
                if para.style and para.style.name and para.style.name.startswith('Heading'):
                    level = para.style.name[-1] if para.style.name[-1].isdigit() else '2'
                    text_parts.append(f"{'#' * (int(level) + 1)} {para.text}")
                else:
                    text_parts.append(para.text)
        
        # Extract tables
        for i, table in enumerate(doc.tables):
            if table.rows:
                text_parts.append(f"\n**Table {i + 1}:**\n")
                
                # Header row
                if table.rows:
                    header_cells = [cell.text.strip() for cell in table.rows[0].cells]
                    text_parts.append("| " + " | ".join(header_cells) + " |")
                    text_parts.append("|" + "|".join(["---"] * len(header_cells)) + "|")
                    
                    # Data rows
                    for row in table.rows[1:]:
                        cells = [cell.text.strip() for cell in row.cells]
                        text_parts.append("| " + " | ".join(cells) + " |")
        
        result = "\n\n".join(text_parts)
        logger.info(f"DOCX processed: {doc_name} - {len(result)} chars extracted")
        return result
        
    except ImportError:
        return f"[DOCX support requires python-docx. Document: {doc_name}]"
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        return f"[DOCX extraction failed for '{doc_name}': {str(e)}]"


def _process_pptx(binary_data: bytes, doc_name: str) -> str:
    """Process PPTX files using python-pptx."""
    try:
        from pptx import Presentation
        
        pptx_file = io.BytesIO(binary_data)
        prs = Presentation(pptx_file)
        
        text_parts = [f"## Presentation: {doc_name}\n"]
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            if slide_texts:
                text_parts.append(f"### Slide {slide_num}\n\n" + "\n\n".join(slide_texts))
        
        result = "\n\n".join(text_parts)
        logger.info(f"PPTX processed: {doc_name} - {len(result)} chars extracted")
        return result
        
    except ImportError:
        return f"[PPTX support requires python-pptx. Document: {doc_name}]"
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        return f"[PPTX extraction failed for '{doc_name}': {str(e)}]"


def _process_xlsx(binary_data: bytes, doc_name: str) -> str:
    """Process XLSX files using openpyxl."""
    try:
        import openpyxl
        
        xlsx_file = io.BytesIO(binary_data)
        wb = openpyxl.load_workbook(xlsx_file, data_only=True)
        
        text_parts = [f"## Spreadsheet: {doc_name}\n"]
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"### Sheet: {sheet_name}\n")
            
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(v.strip() for v in row_values):  # Skip empty rows
                    rows.append(row_values)
            
            if rows:
                # Create markdown table
                if rows:
                    text_parts.append("| " + " | ".join(rows[0]) + " |")
                    text_parts.append("|" + "|".join(["---"] * len(rows[0])) + "|")
                    
                    for row in rows[1:50]:  # Limit to 50 rows
                        text_parts.append("| " + " | ".join(row) + " |")
                    
                    if len(rows) > 50:
                        text_parts.append(f"\n*... ({len(rows) - 50} more rows)*")
        
        result = "\n".join(text_parts)
        logger.info(f"XLSX processed: {doc_name} - {len(result)} chars extracted")
        return result
        
    except ImportError:
        return f"[XLSX support requires openpyxl. Document: {doc_name}]"
    except Exception as e:
        logger.error(f"XLSX extraction error: {e}")
        return f"[XLSX extraction failed for '{doc_name}': {str(e)}]"


def _process_html(binary_data: bytes, doc_name: str) -> str:
    """Process HTML files."""
    from html.parser import HTMLParser
    
    class TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text_parts = []
            self.skip_tags = {'script', 'style', 'head', 'meta', 'link', 'noscript'}
            self.current_skip = False
            
        def handle_starttag(self, tag, attrs):
            if tag in self.skip_tags:
                self.current_skip = True
            elif tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                self.text_parts.append(f"\n{'#' * int(tag[1])} ")
            elif tag in ['p', 'div', 'br']:
                self.text_parts.append("\n")
                
        def handle_endtag(self, tag):
            if tag in self.skip_tags:
                self.current_skip = False
                
        def handle_data(self, data):
            if not self.current_skip:
                text = data.strip()
                if text:
                    self.text_parts.append(text)
    
    try:
        html_content = binary_data.decode('utf-8', errors='replace')
        extractor = TextExtractor()
        extractor.feed(html_content)
        
        content = " ".join(extractor.text_parts)
        return f"## Document: {doc_name}\n\n{content}"
    except Exception as e:
        logger.error(f"HTML extraction error: {e}")
        return f"[HTML extraction failed for '{doc_name}': {str(e)}]"


def _process_csv(binary_data: bytes, doc_name: str) -> str:
    """Process CSV files."""
    import csv
    
    try:
        content = binary_data.decode('utf-8', errors='replace')
        reader = csv.reader(io.StringIO(content))
        
        rows = list(reader)
        if not rows:
            return f"[Empty CSV file: {doc_name}]"
        
        text_parts = [f"## Data: {doc_name}\n"]
        
        # Header
        text_parts.append("| " + " | ".join(rows[0]) + " |")
        text_parts.append("|" + "|".join(["---"] * len(rows[0])) + "|")
        
        # Data rows (limit to 100)
        for row in rows[1:101]:
            text_parts.append("| " + " | ".join(row) + " |")
        
        if len(rows) > 101:
            text_parts.append(f"\n*... ({len(rows) - 101} more rows)*")
        
        return "\n".join(text_parts)
        
    except Exception as e:
        logger.error(f"CSV extraction error: {e}")
        return f"[CSV extraction failed for '{doc_name}': {str(e)}]"


def _process_text(binary_data: bytes, doc_name: str) -> str:
    """Process plain text and markdown files."""
    try:
        content = binary_data.decode('utf-8', errors='replace')
        return f"## Document: {doc_name}\n\n{content}"
    except Exception as e:
        return f"[Text extraction failed for '{doc_name}': {str(e)}]"


def _process_json(binary_data: bytes, doc_name: str) -> str:
    """Process JSON files."""
    import json
    
    try:
        content = binary_data.decode('utf-8', errors='replace')
        data = json.loads(content)
        formatted = json.dumps(data, indent=2)
        return f"## JSON Data: {doc_name}\n\n```json\n{formatted}\n```"
    except Exception as e:
        logger.error(f"JSON processing error: {e}")
        return f"[JSON processing failed for '{doc_name}': {str(e)}]"


def _process_xml(binary_data: bytes, doc_name: str) -> str:
    """Process XML files."""
    try:
        content = binary_data.decode('utf-8', errors='replace')
        return f"## XML Document: {doc_name}\n\n```xml\n{content}\n```"
    except Exception as e:
        return f"[XML extraction failed for '{doc_name}': {str(e)}]"
